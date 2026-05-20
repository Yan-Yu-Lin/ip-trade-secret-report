#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.11"
# dependencies = ["playwright", "python-pptx", "Pillow"]
# ///

"""
Screenshot each slide of the HTML deck and assemble into a PPTX.
Usage: uv run capture_and_build.py
"""

import asyncio
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from playwright.async_api import async_playwright

HTML_PATH = Path(__file__).parent.parent / "ppt" / "index.html"
NARRATIVE_PATH = Path(__file__).parent.parent / "08-完整連貫敘述-v3.md"
OUTPUT_DIR = Path(__file__).parent
SCREENSHOTS_DIR = OUTPUT_DIR / "slides"
OUTPUT_PPTX = OUTPUT_DIR / "presentation.pptx"

SLIDE_WIDTH = 1920
SLIDE_HEIGHT = 1080
TOTAL_SLIDES = 24


def parse_speaker_notes(narrative_path: Path) -> dict[int, str]:
    """Parse narrative markdown and roughly map sections to slide numbers."""
    text = narrative_path.read_text(encoding="utf-8")
    sections = re.split(r'\n## ', text)

    # Rough mapping: narrative sections to slide ranges
    mapping = {}
    slide_section_map = {
        1: "從一個基本問題開始",
        2: "專利是什麼",
        3: "專利是什麼",
        4: "那營業秘密是什麼",
        5: "那為什麼不全部用專利",
        6: "那營業秘密怎麼拿來賺錢",
        7: "那營業秘密怎麼拿來賺錢",
        8: "但如果要授權",
        9: "但如果要授權",
        10: "但如果要授權",
        11: "第二種：打包授權",
        12: "第二種：打包授權",
        13: "第二種：打包授權",
        14: "第二種：打包授權",
        15: "第一種：自己用",
        16: "第一種：自己用",
        17: "第一種：自己用",
        18: "第一種：自己用",
        19: "第三種：保護不好",
        20: "第三種：保護不好",
        21: "第三種：保護不好",
        22: "第三種：保護不好",
        23: "營業秘密怎麼保護",
        24: "把所有東西串起來",
    }

    for slide_num, keyword in slide_section_map.items():
        for section in sections:
            if keyword in section:
                # Take first 500 chars as notes
                clean = section.strip()[:800]
                mapping[slide_num] = clean
                break

    return mapping


async def capture_slides():
    """Use Playwright to screenshot each slide."""
    SCREENSHOTS_DIR.mkdir(exist_ok=True)

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(viewport={"width": SLIDE_WIDTH, "height": SLIDE_HEIGHT})

        file_url = f"file://{HTML_PATH.resolve()}"
        print(f"Loading {file_url}")
        await page.goto(file_url, wait_until="networkidle", timeout=60000)

        # Wait for WebGL and animations to settle
        await page.wait_for_timeout(3000)

        # Force static/low-power mode for clean screenshots
        await page.evaluate("window.__setLowPowerMode && window.__setLowPowerMode(true, {persist:false})")
        await page.wait_for_timeout(500)

        # Get total slides
        total = await page.evaluate("document.querySelectorAll('.slide').length")
        print(f"Found {total} slides")

        for i in range(total):
            # Navigate to slide
            await page.evaluate(f"""
                const deck = document.getElementById('deck');
                deck.style.transition = 'none';
                deck.style.transform = 'translateX(-{i * 100}vw)';
            """)
            await page.wait_for_timeout(300)

            # Force all animations to final state on this slide
            await page.evaluate(f"""
                const slides = document.querySelectorAll('.slide');
                const slide = slides[{i}];
                slide.querySelectorAll('[data-anim]').forEach(el => {{
                    el.style.opacity = '1';
                    el.style.transform = 'none';
                }});
            """)
            await page.wait_for_timeout(200)

            path = SCREENSHOTS_DIR / f"slide-{i+1:02d}.png"
            await page.screenshot(path=str(path))
            print(f"  Captured slide {i+1}/{total}")

        await browser.close()


def build_pptx(notes_map: dict[int, str]):
    """Assemble screenshots into PPTX with speaker notes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for i in range(TOTAL_SLIDES):
        img_path = SCREENSHOTS_DIR / f"slide-{i+1:02d}.png"
        if not img_path.exists():
            print(f"  Warning: {img_path} not found, skipping")
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Add screenshot as full-bleed background
        slide.shapes.add_picture(
            str(img_path),
            Emu(0), Emu(0),
            prs.slide_width, prs.slide_height
        )

        # Add speaker notes
        if (i + 1) in notes_map:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_map[i + 1]

        print(f"  Added slide {i+1}")

    prs.save(str(OUTPUT_PPTX))
    print(f"\nSaved to {OUTPUT_PPTX}")


async def main():
    print("=== Step 1: Capturing slides ===")
    await capture_slides()

    print("\n=== Step 2: Parsing speaker notes ===")
    notes_map = parse_speaker_notes(NARRATIVE_PATH)
    print(f"  Mapped notes for {len(notes_map)} slides")

    print("\n=== Step 3: Building PPTX ===")
    build_pptx(notes_map)

    print("\nDone!")


if __name__ == "__main__":
    asyncio.run(main())
